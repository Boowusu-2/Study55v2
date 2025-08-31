import os
import traceback
import sys
import base64
from typing import Optional

def read_txt(path: str) -> str:
    """Extract text from TXT files"""
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
            print(f"Successfully read TXT file: {len(content)} characters", file=sys.stderr)
            return content
    except Exception as e:
        print(f"TXT_READ_ERROR for {path}: {e}", file=sys.stderr)
        return ""

def read_pdf(path: str) -> str:
    """Extract text from PDF files"""
    try:
        import PyPDF2
        text_parts = []
        with open(path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            print(f"PDF has {len(reader.pages)} pages", file=sys.stderr)
            for i, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text() or ""
                    text_parts.append(page_text)
                    print(f"Page {i+1}: extracted {len(page_text)} characters", file=sys.stderr)
                except Exception as page_error:
                    print(f"Error reading page {i+1}: {page_error}", file=sys.stderr)
        
        result = "\n".join(text_parts)
        print(f"Total PDF text: {len(result)} characters", file=sys.stderr)
        return result
    except ImportError:
        print("PyPDF2 not installed. Install with: pip install PyPDF2", file=sys.stderr)
        return ""
    except Exception as e:
        print(f"PDF_PARSE_ERROR for {path}: {e}", file=sys.stderr)
        print(traceback.format_exc(), file=sys.stderr)
        return ""

def read_docx(path: str) -> str:
    """Extract text from DOCX files"""
    try:
        import docx
        doc = docx.Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        result = "\n".join(paragraphs)
        print(f"DOCX extracted {len(paragraphs)} paragraphs, {len(result)} characters", file=sys.stderr)
        return result
    except ImportError:
        print("python-docx not installed. Install with: pip install python-docx", file=sys.stderr)
        return ""
    except Exception as e:
        print(f"DOCX_PARSE_ERROR for {path}: {e}", file=sys.stderr)
        print(traceback.format_exc(), file=sys.stderr)
        return ""

def read_doc(path: str) -> str:
    """Extract text from DOC files"""
    try:
        # Try using python-docx2txt for .doc files
        import docx2txt
        result = docx2txt.process(path)
        print(f"DOC extracted {len(result)} characters", file=sys.stderr)
        return result or ""
    except ImportError:
        print("docx2txt not installed. Install with: pip install docx2txt", file=sys.stderr)
        try:
            # Fallback: try to read as plain text (might work for some .doc files)
            return read_txt(path)
        except:
            return ""
    except Exception as e:
        print(f"DOC_PARSE_ERROR for {path}: {e}", file=sys.stderr)
        return ""

def read_pptx(path: str) -> str:
    """Extract text from PPTX files"""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text_parts = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            if slide_text:
                text_parts.append(f"Slide {i+1}:\n" + "\n".join(slide_text))
        
        result = "\n\n".join(text_parts)
        print(f"PPTX extracted {len(prs.slides)} slides, {len(result)} characters", file=sys.stderr)
        return result
    except ImportError:
        print("python-pptx not installed. Install with: pip install python-pptx", file=sys.stderr)
        return ""
    except Exception as e:
        print(f"PPTX_PARSE_ERROR for {path}: {e}", file=sys.stderr)
        return ""

def read_image_with_ocr(path: str) -> str:
    """Extract text from images using OCR (Optical Character Recognition)"""
    try:
        # Try using EasyOCR first (better for handwritten text)
        try:
            import easyocr
            print(f"Using EasyOCR for image: {path}", file=sys.stderr)
            
            # Initialize EasyOCR reader (downloads models on first use)
            reader = easyocr.Reader(['en'], gpu=False)  # Use CPU for compatibility
            
            # Read the image
            results = reader.readtext(path)
            
            # Extract text from results
            text_parts = []
            for (bbox, text, confidence) in results:
                if confidence > 0.5:  # Only include text with confidence > 50%
                    text_parts.append(text)
                    print(f"OCR detected: '{text}' (confidence: {confidence:.2f})", file=sys.stderr)
            
            result = "\n".join(text_parts)
            print(f"EasyOCR extracted {len(text_parts)} text blocks, {len(result)} characters", file=sys.stderr)
            return result
            
        except ImportError:
            print("EasyOCR not installed. Trying Tesseract...", file=sys.stderr)
            
            # Fallback to Tesseract OCR
            try:
                import pytesseract
                from PIL import Image
                
                # Open the image
                image = Image.open(path)
                
                # Extract text using Tesseract
                text = pytesseract.image_to_string(image)
                
                print(f"Tesseract extracted {len(text)} characters", file=sys.stderr)
                return text.strip()
                
            except ImportError:
                print("Tesseract not available. Install with: pip install pytesseract pillow", file=sys.stderr)
                return ""
            except Exception as tesseract_error:
                print(f"Tesseract OCR error: {tesseract_error}", file=sys.stderr)
                return ""
                
    except Exception as e:
        print(f"OCR_ERROR for {path}: {e}", file=sys.stderr)
        print(traceback.format_exc(), file=sys.stderr)
        return ""

def read_image_with_google_vision(path: str) -> str:
    """Extract text from images using Google Cloud Vision API (if available)"""
    try:
        from google.cloud import vision
        
        # Initialize the client
        client = vision.ImageAnnotatorClient()
        
        # Read the image file
        with open(path, 'rb') as image_file:
            content = image_file.read()
        
        # Create image object
        image = vision.Image(content=content)
        
        # Perform text detection
        response = client.text_detection(image=image)
        texts = response.text_annotations
        
        if texts:
            # The first element contains the entire text
            full_text = texts[0].description
            print(f"Google Vision extracted {len(full_text)} characters", file=sys.stderr)
            return full_text
        else:
            print("No text detected in image", file=sys.stderr)
            return ""
            
    except ImportError:
        print("Google Cloud Vision not installed. Install with: pip install google-cloud-vision", file=sys.stderr)
        return ""
    except Exception as e:
        print(f"Google Vision OCR error: {e}", file=sys.stderr)
        return ""

def extract_text_from_image(path: str) -> str:
    """Extract text from image files using multiple OCR methods"""
    print(f"Processing image file: {path}", file=sys.stderr)
    
    # Try Google Cloud Vision first (if available and configured)
    text = read_image_with_google_vision(path)
    if text and text.strip():
        return text
    
    # Fallback to local OCR methods
    text = read_image_with_ocr(path)
    if text and text.strip():
        return text
    
    print(f"No text could be extracted from image: {path}", file=sys.stderr)
    return ""

def extract_text_from_files(file_paths: list) -> str:
    """
    Extract text from multiple files and return combined text.
    
    Args:
        file_paths: List of file paths to process
        
    Returns:
        Combined text from all files
    """
    print(f"Processing {len(file_paths)} files", file=sys.stderr)
    
    outputs = []
    for file_path in file_paths:
        print(f"Processing file: {file_path}", file=sys.stderr)
        
        if not os.path.exists(file_path):
            print(f"File does not exist: {file_path}", file=sys.stderr)
            continue
            
        file_size = os.path.getsize(file_path)
        print(f"File size: {file_size} bytes", file=sys.stderr)
        
        ext = os.path.splitext(file_path)[1].lower()
        print(f"File extension: {ext}", file=sys.stderr)
        
        text = ""
        if ext == ".txt":
            text = read_txt(file_path)
        elif ext == ".pdf":
            text = read_pdf(file_path)
        elif ext == ".docx":
            text = read_docx(file_path)
        elif ext == ".doc":
            text = read_doc(file_path)
        elif ext in [".ppt", ".pptx"]:
            text = read_pptx(file_path)
        elif ext.lower() in [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif", ".gif", ".webp"]:
            text = extract_text_from_image(file_path)
        else:
            print(f"Unsupported file type: {ext}", file=sys.stderr)
            continue
        
        if text and text.strip():
            filename = os.path.basename(file_path)
            outputs.append(f"\n\n=== Content from {filename} ===\n{text}")
            print(f"Added content from {filename}: {len(text)} characters", file=sys.stderr)
        else:
            print(f"No text extracted from {file_path}", file=sys.stderr)

    final_output = "".join(outputs)
    print(f"Final output length: {len(final_output)} characters", file=sys.stderr)
    
    return final_output.strip()
