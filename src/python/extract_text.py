#!/usr/bin/env python3
import sys
import os
import traceback

def read_txt(path: str) -> str:
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
            print(f"Successfully read TXT file: {len(content)} characters", file=sys.stderr)
            return content
    except Exception as e:
        print(f"TXT_READ_ERROR for {path}: {e}", file=sys.stderr)
        return ""

def read_pdf(path: str) -> str:
    try:
        import PyPDF2
        text_parts = []
        with open(path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            print(f"PDF has {len(reader.pages)} pages", file=sys.stderr)
            for i, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text() or ""
                    text_parts.append(page_text)
                    print(f"Page {i+1}: extracted {len(page_text)} characters", file=sys.stderr)
                except Exception as page_error:
                    print(f"Error reading page {i+1}: {page_error}", file=sys.stderr)
        
        result = "\n".join(text_parts)
        print(f"Total PDF text: {len(result)} characters", file=sys.stderr)
        return result
    except ImportError:
        print("PyPDF2 not installed. Install with: pip install PyPDF2", file=sys.stderr)
        return ""
    except Exception as e:
        print(f"PDF_PARSE_ERROR for {path}: {e}", file=sys.stderr)
        print(traceback.format_exc(), file=sys.stderr)
        return ""

def read_docx(path: str) -> str:
    try:
        import docx
        doc = docx.Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        result = "\n".join(paragraphs)
        print(f"DOCX extracted {len(paragraphs)} paragraphs, {len(result)} characters", file=sys.stderr)
        return result
    except ImportError:
        print("python-docx not installed. Install with: pip install python-docx", file=sys.stderr)
        return ""
    except Exception as e:
        print(f"DOCX_PARSE_ERROR for {path}: {e}", file=sys.stderr)
        print(traceback.format_exc(), file=sys.stderr)
        return ""

def read_doc(path: str) -> str:
    try:
        # Try using python-docx2txt for .doc files
        import docx2txt
        result = docx2txt.process(path)
        print(f"DOC extracted {len(result)} characters", file=sys.stderr)
        return result or ""
    except ImportError:
        print("docx2txt not installed. Install with: pip install docx2txt", file=sys.stderr)
        try:
            # Fallback: try to read as plain text (might work for some .doc files)
            return read_txt(path)
        except:
            return ""
    except Exception as e:
        print(f"DOC_PARSE_ERROR for {path}: {e}", file=sys.stderr)
        return ""

def read_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text_parts = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            if slide_text:
                text_parts.append(f"Slide {i+1}:\n" + "\n".join(slide_text))
        
        result = "\n\n".join(text_parts)
        print(f"PPTX extracted {len(prs.slides)} slides, {len(result)} characters", file=sys.stderr)
        return result
    except ImportError:
        print("python-pptx not installed. Install with: pip install python-pptx", file=sys.stderr)
        return ""
    except Exception as e:
        print(f"PPTX_PARSE_ERROR for {path}: {e}", file=sys.stderr)
        return ""

def main():
    print(f"Python script started with {len(sys.argv)} arguments", file=sys.stderr)
    print(f"Arguments: {sys.argv}", file=sys.stderr)
    
    if len(sys.argv) < 2:
        print("No file paths provided", file=sys.stderr)
        print("")  # Empty output
        return

    outputs = []
    for file_path in sys.argv[1:]:
        print(f"Processing file: {file_path}", file=sys.stderr)
        
        if not os.path.exists(file_path):
            print(f"File does not exist: {file_path}", file=sys.stderr)
            continue
            
        file_size = os.path.getsize(file_path)
        print(f"File size: {file_size} bytes", file=sys.stderr)
        
        ext = os.path.splitext(file_path)[1].lower()
        print(f"File extension: {ext}", file=sys.stderr)
        
        text = ""
        if ext == ".txt":
            text = read_txt(file_path)
        elif ext == ".pdf":
            text = read_pdf(file_path)
        elif ext == ".docx":
            text = read_docx(file_path)
        elif ext == ".doc":
            text = read_doc(file_path)
        elif ext in [".ppt", ".pptx"]:
            text = read_pptx(file_path)
        else:
            print(f"Unsupported file type: {ext}", file=sys.stderr)
            continue
        
        if text and text.strip():
            filename = os.path.basename(file_path)
            outputs.append(f"\n\n=== Content from {filename} ===\n{text}")
            print(f"Added content from {filename}: {len(text)} characters", file=sys.stderr)
        else:
            print(f"No text extracted from {file_path}", file=sys.stderr)

    final_output = "".join(outputs)
    print(f"Final output length: {len(final_output)} characters", file=sys.stderr)
    
    if final_output.strip():
        print(final_output)  # This goes to stdout
    else:
        print("No content extracted from any files", file=sys.stderr)
        print("")  # Empty output to stdout

if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        print(f"Script error: {e}", file=sys.stderr)
        print(traceback.format_exc(), file=sys.stderr)
        print("")  # Empty output